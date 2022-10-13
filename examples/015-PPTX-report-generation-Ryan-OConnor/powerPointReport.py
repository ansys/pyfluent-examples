"""
015-PPTX-report-generation
============================
These examples show you how you can use Fluent capabilities from Python to perform
Fluent simulations. This includes geometry import, Fluent's meshing workflows,
setting up and running the solver, and reviewing the results using Fluent's
postprocessing capabilities.
"""

# Create PowerPoint Report with Fluent Results
# This script automatically generates a PPTX report based on Fluent results.
# The PPTX contains a title slide, a table of computed report definitions,
# images of user created graphics objects
# (mesh, contours, vectors, pathlines, and particle tracks), and
# images of charts (residuals and report plots).
# It can be easily modified to work with any PowerPoint template.

import ansys.fluent.core as pyfluent

# Impoer modules
from pptx import Presentation
from pptx.util import Inches

# User Inputs
case_filename = "elbow.cas.h5"
data_filename = "elbow.dat.h5"
template_filename = "templatePyAnsys.pptx"

# Determine Placeholder Index for Given PPTX Template
# Analyze_ppt is taken from
# https://github.com/chris1610/pbpython/blob/master/code/analyze_ppt.pyhttps:
# //github.com/chris1610/pbpython/blob/master/code/analyze_ppt.py
# * Analyze_ppt is used to determine the placeholder indices in the PPTX template.
# It generates an annotated PPTX named labelled_template.pptx.
# This function is only required if the placeholder indices are not known.
# * Review the generated labelled_template.pptx and
# determine slide desired layout (indices start at 0) and
# for each layout the required placeholders
# * Images can be inserted using placeholders of type:Picture.
# Alternatively images can be inserted anywhere in the slide
# based on location (not covered in this script)
# * If template does not include the desired layout,
# edit template first and add layout with desired placeholder arrangement


def analyze_ppt(input, output):
    """Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.
    """
    prs = Presentation(input)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = "Title for Layout {}".format(index)
        except AttributeError:
            continue
            # print("No Title for Layout {}".format(index))
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if "Title" not in shape.text:
                        shape.text = "Placeholder index:{} type:{}".format(
                            phf.idx, shape.name
                        )
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                # print('{} {}'.format(phf.idx, shape.name))
    prs.save(output)


analyze_ppt(template_filename, "labelled_template.pptx")

# Load Fluent Model
session = pyfluent.launch_fluent()

# Check server status
session.check_health()

# Read case file and data
session.solver.tui.file.read_case(case_filename)
session.solver.tui.file.read_data(data_filename)

root = session.solver.root

# Open PPTX Template
prs = Presentation(template_filename)

# Add Title Slide

# slide layout 0 is selected based on particulate template
slide = prs.slides.add_slide(prs.slide_layouts[0])

# placeholder 10 holds the title text for particular template
title = slide.placeholders[10]
title.text = case_filename

# placeholder 12 holds the subtitle text for particulare template
subtitle = slide.placeholders[12]
subtitle.text = "CFD Simulation Results"

# Add Report Definition Slide
# Collect Report Definition Data
repdef = []
repcalc = []
reportList = root.solution.report_definitions.surface.get_object_names()

for report in reportList:
    for key, value in root.solution.report_definitions.compute(report_defs=[report])[
        0
    ].items():
        repdef.append(key)
        repcalc.append(value[0])

if reportList:  # if there are reports
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    title = slide.shapes.title
    title.text = "Report Definitions"

    # Create a table of appropriate size
    x, y, cx, cy = Inches(2), Inches(2), Inches(8), Inches(1.5)
    shape = slide.shapes.add_table(len(repdef) + 1, 2, x, y, cx, cy)
    table = shape.table

    # Create table headers
    cell = table.cell(0, 0)
    cell.text = "Report Definition"
    cell = table.cell(0, 1)
    cell.text = "Value"

    # Populate table with Report Definition Data
    for index, value in enumerate(repdef):
        cell = table.cell(index + 1, 0)
        cell.text = value
        cell = table.cell(index + 1, 1)
        cell.text = format(repcalc[index], ".4E")

# Function that will be used later to resize pictures to fit their placeholder
def adjust_picture_to_fit(picture):
    available_width = picture.width
    available_height = picture.height
    image_width, image_height = picture.image.size
    placeholder_aspect_ratio = float(available_width) / float(available_height)
    image_aspect_ratio = float(image_width) / float(image_height)
    pos_left, pos_top = picture.left, picture.top

    picture.crop_top = 0
    picture.crop_left = 0
    picture.crop_bottom = 0
    picture.crop_right = 0

    # ---if the placeholder is "wider" in aspect, shrink the picture width while
    # ---maintaining the image aspect ratio
    if placeholder_aspect_ratio > image_aspect_ratio:
        picture.height = available_height
        picture.width = int(image_aspect_ratio * available_height)
    # ---otherwise shrink the height
    else:
        picture.width = available_width
        picture.height = int(available_width / image_aspect_ratio)
    picture.left, picture.top = pos_left, pos_top


# Add Graphics Slides
# Dictionary corresponding to graphics objects.
# Once scenes are enabled with the Settings API they can be added here.
Images = {
    "mesh": root.results.graphics.mesh.get_object_names(),
    "contour": root.results.graphics.contour.get_object_names(),
    "vector": root.results.graphics.vector.get_object_names(),
    "pathlines": root.results.graphics.pathline.get_object_names(),
    "particletracks": root.results.graphics.particle_track.get_object_names(),
}

session.solver.tui.display.set.picture.use_window_resolution("no")

# Loop through the objects of each key in the Images dictionary,
# save picture, and insert picture into PPTX
for key, value in Images.items():
    for image in value:
        graph_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(graph_slide_layout)
        title = slide.shapes.title
        title.text = image
        session.solver.tui.display.save_picture(image + ".png")
        placeholder = slide.placeholders[14]
        pic = placeholder.insert_picture(image + ".png")
        adjust_picture_to_fit(pic)

# Add residual plot.
session.solver.tui.plot.residuals()
session.solver.tui.display.save_picture("residuals.png")

graph_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(graph_slide_layout)
title = slide.shapes.title
title.text = "Residuals"

placeholder = slide.placeholders[14]
pic = placeholder.insert_picture("residuals.png")
adjust_picture_to_fit(pic)

# Add Report Charts
session.solver.tui.solve.report_plots.list()
# Add charts for each individual report plot.

# Generate a list of all report plots.
# Using Scheme Eval until Settings API can be used.
reportdefs_string = session.scheme_eval.exec(
    ('(ti-menu-load-string "solve/report-plots/list ")',)
)
reportdefs = reportdefs_string.split()
reportdefs.pop(0)

# Loop through all report plots,
# generate a slide titled with report plot name,
# including image of plot
if reportdefs[0] == "invalid":
    print("There are no report definitions.")
else:
    for report in reportdefs:
        session.solver.tui.solve.report_plots.plot(report)
        session.solver.tui.display.save_picture(report + ".png")

        graph_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(graph_slide_layout)
        title = slide.shapes.title
        title.text = report

        placeholder = slide.placeholders[14]
        pic = placeholder.insert_picture(report + ".png")
        adjust_picture_to_fit(pic)

# Save and close PPTX
prs.save(case_filename + "_report.pptx")

# End current session
# session.exit()
