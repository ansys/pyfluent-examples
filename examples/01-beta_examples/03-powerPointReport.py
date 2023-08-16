"""

##########################################################
PPTX Report Generation
##########################################################
.. _Write example object:

Objective
===============

This example takes an existing Fluent .cas file and goes through the steps necessary 
to post process the results and extract the information to a PPTX report. In this 
example we are using Fluent's inbuild postprocessing capabilities. 

The steps include:

* This script automatically generates a PPTX report based on Fluent results.
* The PPTX contains a title slide, a table of computed report definitions,
* images of user created graphics objects
* (mesh, contours, vectors, pathlines, and particle tracks), and
* images of charts (residuals and report plots).
* It can be easily modified to work with any PowerPoint template.

.. image:: ../../_static/PPTXReport_mesh.png
   :align: center
   :alt: Mesh

"""
####################################################################################
# Import required libraries/modules
# ==================================================================================

from pathlib import Path

import ansys.fluent.core as pyfluent
from ansys.fluent.core import examples

from pptx import Presentation
from pptx.util import Inches

####################################################################################
# User Inputs
# ==================================================================================
# In this example the report is generated using the elbow case taken from the Fluent Tutorial.
# The steps needed to generate a pptx report could be applied to any Fluent simulation.
case_filename = "elbow.cas.h5"
data_filename = "elbow.dat.h5"
template_filename = "PPTX-report-generation/templatePyAnsys.pptx"

import_case_filename = examples.download_file(
    "elbow.cas.h5",
    "pyfluent/examples/PPTX-report-generation",
)

import_data_filename = examples.download_file(
    "elbow.dat.h5",
    "pyfluent/examples/PPTX-report-generation",
)

import_template_filename = examples.download_file(
    "templatePyAnsys.pptx",
    "pyfluent/examples/PPTX-report-generation",
)

####################################################################################
# Organize PPTX Template
# ==================================================================================
# Determine Placeholder Index for Given PPTX Template
# Analyze_ppt is taken from
# https://github.com/chris1610/pbpython/blob/master/code/analyze_ppt.pyhttps:
# //github.com/chris1610/pbpython/blob/master/code/analyze_ppt.py
#
# * Analyze_ppt is used to determine the placeholder indices in the PPTX template.
# * It generates an annotated PPTX named labelled_template.pptx.
# * This function is only required if the placeholder indices are not known.
# * Review the generated labelled_template.pptx and determine slide desired layout and placeholders
# * Images can be inserted using placeholders of type:Picture.
# * Alternatively images can be inserted anywhere in the slide based on location (not covered in this script)
# * If template does not include the desired layout, edit template first and add layout with desired placeholder arrangement


def analyze_ppt(input, output):
    """Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.
    """
    prs = Presentation(input)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = "Title for Layout {}".format(index)
        except AttributeError:
            continue
            # print("No Title for Layout {}".format(index))
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if "Title" not in shape.text:
                        shape.text = "Placeholder index:{} type:{}".format(
                            phf.idx, shape.name
                        )
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                # print('{} {}'.format(phf.idx, shape.name))
    prs.save(output)


analyze_ppt(
    import_template_filename,
    str(Path(pyfluent.EXAMPLES_PATH) / "labelled_template.pptx"),
)


####################################################################################
# Open Fluent
# ==================================================================================
session = pyfluent.launch_fluent(product_version = "23.1.0")

####################################################################################
# Read case file and data
# ==================================================================================
session.tui.file.read_case(import_case_filename)
session.tui.file.read_data(import_data_filename)

####################################################################################
# Initialize and run case
# ==================================================================================
# report data not stored in saved .dat files so need to run before creating plots.
session.solution.initialization.hybrid_initialize()
session.tui.display.set.windows.aspect_ratio(1920,1080)
session.tui.solve.set.number_of_iterations(50)
session.tui.solve.iterate()

####################################################################################
# Open PPTX Template
# ==================================================================================
prs = Presentation(import_template_filename)

####################################################################################
# Add Title Slide
# ==================================================================================
#
# * slide layout 0 is selected based on particulate template
# * placeholder 10 holds the title text for particular template
# * placeholder 12 holds the subtitle text for particulare template
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.placeholders[10]
title.text = case_filename
subtitle = slide.placeholders[12]
subtitle.text = "CFD Simulation Results"

####################################################################################
# Add Report Definition
# ==================================================================================
# Look at list of reports in the Fluent simulation
# Collect Report Definition Data and assign to arrays
repdef = []
repcalc = []
reportList = session.solution.report_definitions.surface.get_object_names()

for report in reportList:
    dataList = session.solution.report_definitions.compute(report_defs=[report])
    keyList = list(dataList[0].keys())
    valueList = list(dataList[0].values())
    for key in keyList:
        repdef.append(key)
    for value in valueList:           
        repcalc.append(value[0])

####################################################################################
# Create report slide(s)
# ==================================================================================
# If there are reports. For each report:
#
# * Add a slide (layout option 3)
# * Add a table with correct numner of columns and rows
# * Add table headers
# * Copy the report names and values
#
# .. image:: ../../_static/PPTXReport_Table.png
#    :align: center
#    :alt: Table


if reportList: 
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    title = slide.shapes.title
    title.text = "Report Definitions"

    x, y, cx, cy = Inches(2), Inches(2), Inches(8), Inches(1.5)
    shape = slide.shapes.add_table(len(repdef) + 1, 2, x, y, cx, cy)
    table = shape.table

    cell = table.cell(0, 0)
    cell.text = "Report Definition"
    cell = table.cell(0, 1)
    cell.text = "Value"

    for index, value in enumerate(repdef):
        cell = table.cell(index + 1, 0)
        cell.text = value
        cell = table.cell(index + 1, 1)
        cell.text = format(repcalc[index], ".4E")

####################################################################################
# Define adjust_picture_to_fit function
# ==================================================================================
# Function that will be used later to resize pictures to fit their placeholder
# It will apply the following modifications:
#
# * Do not crop the image.
# * If the image is "taller" in aspect than the available height,
# * shrink the picture to use the allowable height
# * Otherwise scale the picture to use the full 
# * Always maintain the image aspect ratio


def adjust_picture_to_fit(picture):
    available_width = picture.width
    available_height = picture.height
    image_width, image_height = picture.image.size
    placeholder_aspect_ratio = float(available_width) / float(available_height)
    image_aspect_ratio = float(image_width) / float(image_height)
    pos_left, pos_top = picture.left, picture.top

    picture.crop_top = 0
    picture.crop_left = 0
    picture.crop_bottom = 0
    picture.crop_right = 0

    if placeholder_aspect_ratio > image_aspect_ratio:
        picture.height = available_height
        picture.width = int(image_aspect_ratio * available_height)
    else:
        picture.width = available_width
        picture.height = int(available_width / image_aspect_ratio)
    picture.left, picture.top = pos_left, pos_top

####################################################################################
# Add Graphics Slides
# ==================================================================================
# Create dictionary corresponding to graphics objects defined in the Fluent simulation
# Once scenes are enabled with the Settings API they can be added here.
# Loop through the objects of each key in the Images dictionary and:
#
# * Add a slide into the PPTX (layout[6])
# * Add slide title
# * dispay the image in Fluent session.tui.display.set.picture.use_window_resolution("no")
# * save picture
# * insert picture into PPTX
# * use adjust_picture_to_fit function to adjust size
#
# .. image:: ../../_static/PPTXReport_velocity.png
#    :align: center
#    :alt: Velocity Magnitude
#
# .. image:: ../../_static/PPTXReport_temperature.png
#    :align: center
#    :alt: Temperature
Images = {
    "mesh": session.results.graphics.mesh.get_object_names(),
    "contour": session.results.graphics.contour.get_object_names(),
    "vector": session.results.graphics.vector.get_object_names(),
    "pathlines": session.results.graphics.pathline.get_object_names(),
    "particletracks": session.results.graphics.particle_track.get_object_names(),
}
for key, value in Images.items():
    for image in value:
        graph_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(graph_slide_layout)
        title = slide.shapes.title
        title.text = image
        session.tui.display.objects.display(image)
        img_path = str((Path(pyfluent.EXAMPLES_PATH) / str(image)).with_suffix(".png"))
        session.tui.display.set.picture.driver.png()
        session.tui.display.set.picture.use_window_resolution("no")
        session.tui.display.set.picture.x_resolution(1920)
        session.tui.display.set.picture.y_resolution(1080)
        session.tui.display.save_picture(img_path)
        placeholder = slide.placeholders[14]
        pic = placeholder.insert_picture(img_path)
        adjust_picture_to_fit(pic)

####################################################################################
# Add residual plot
# ==================================================================================
#
# * Export residual plot from fluent
# * Add slide (layout[6]) to pptx
# * Adjust picture size and add to pptx
session.tui.plot.residuals()
session.tui.display.save_picture(str(Path(pyfluent.EXAMPLES_PATH) / "residuals.png"))

graph_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(graph_slide_layout)
title = slide.shapes.title
title.text = "Residuals"

placeholder = slide.placeholders[14]
pic = placeholder.insert_picture(str(Path(pyfluent.EXAMPLES_PATH) / "residuals.png"))
adjust_picture_to_fit(pic)

####################################################################################
# Add Report Charts
# ==================================================================================
#
# Add charts for each individual report plot:
#
# * Generate a list of all report plots. Take string provided by Fluent and divide into individual plots.
# * Using Scheme Eval until Settings API can be used.
#
# If there are report plots present, loop over reports. For each:
#
# * Plot report 
# * Export .png image
# * Add slide to pptx
# * Add title with report plot name
# * Add the .png
# * Adjust to fit
#
# .. image:: ../../_static/PPTXReport_OutletT.png
#    :align: center
#    :alt: OutletTemperature
session.tui.solve.report_plots.list()
reportdefs_string = session.scheme_eval.exec(
    ('(ti-menu-load-string "solve/report-plots/list ")',)
)
reportdefs = reportdefs_string.split()
reportdefs.pop(0)
if reportdefs[0] == "invalid":
    print("There are no report definitions.")
else:
    for report in reportdefs:
        session.tui.solve.report_plots.plot(report)
        session.tui.display.save_picture(
            str((Path(pyfluent.EXAMPLES_PATH) / str(report)).with_suffix(".png"))
        )

        graph_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(graph_slide_layout)
        title = slide.shapes.title
        title.text = report

        placeholder = slide.placeholders[14]
        pic = placeholder.insert_picture(
            str((Path(pyfluent.EXAMPLES_PATH) / str(report)).with_suffix(".png"))
        )
        adjust_picture_to_fit(pic)

####################################################################################
# Save and close PPTX
# ==================================================================================
prs.save(str((Path(pyfluent.EXAMPLES_PATH) / case_filename).with_suffix(".pptx")))
