"""
016-Part-1-Create-PPTX
===================
These examples show you how you can generate PPTX and use Fluent capabilities
from Python to perform Fluent simulations. This includes geometry import,
Fluent's meshing workflows, setting up and running the solver,
and reviewing the results using Fluent's postprocessing capabilities.
"""

# import modules
from pptx import Presentation
from pptx.util import Inches


# Generate PPTX
def generatePPTX(session):
    template_filename = "template.pptx"

    prs = Presentation(template_filename)

    # slide layout 0 is selected based on particulate template
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # placeholder 10 holds the title text for particular template
    title = slide.placeholders[10]
    title.text = "Pipe Elbow Report"

    # placeholder 12 holds the subtitle text for particular template
    subtitle = slide.placeholders[12]
    subtitle.text = "Automatically Generated Report"

    # Collect Report Definition Data
    repdef = []
    repcalc = []
    reportList = session.solution.report_definitions.surface.get_object_names()
    for report in reportList:
        for key, value in session.solution.report_definitions.compute(
            report_defs=[report]
        )[0].items():
            repdef.append(key)
            repcalc.append(value[0])

    if reportList:  # if there are reports
        slide = prs.slides.add_slide(prs.slide_layouts[3])
        title = slide.shapes.title
        title.text = "Report Definitions"

        # Create a table of appropriate size
        x, y, cx, cy = Inches(2), Inches(2), Inches(8), Inches(1.5)
        shape = slide.shapes.add_table(len(repdef) + 1, 2, x, y, cx, cy)
        table = shape.table

        # Create table headers
        cell = table.cell(0, 0)
        cell.text = "Report Definition"
        cell = table.cell(0, 1)
        cell.text = "Value"

        # Populate table with Report Definition Data
        for index, value in enumerate(repdef):
            cell = table.cell(index + 1, 0)
            cell.text = value
            cell = table.cell(index + 1, 1)
            cell.text = format(repcalc[index], ".4E")

    # Function that will be used later to resize pictures to fit their placeholder
    def adjust_picture_to_fit(picture):
        available_width = picture.width
        available_height = picture.height
        image_width, image_height = picture.image.size
        placeholder_aspect_ratio = float(available_width) / float(available_height)
        image_aspect_ratio = float(image_width) / float(image_height)
        pos_left, pos_top = picture.left, picture.top

        picture.crop_top = 0
        picture.crop_left = 0
        picture.crop_bottom = 0
        picture.crop_right = 0

        # ---if the placeholder is "wider" in aspect, shrink the picture width while
        # ---maintaining the image aspect ratio
        if placeholder_aspect_ratio > image_aspect_ratio:
            picture.height = available_height
            picture.width = int(image_aspect_ratio * available_height)
        # ---otherwise shrink the height
        else:
            picture.width = available_width
            picture.height = int(available_width / image_aspect_ratio)
        picture.left, picture.top = pos_left, pos_top

    # Dictionary corresponding to graphics objects.
    # Once scenes are enabled with the Settings API they can be added here.
    Images = {
        "mesh": session.results.graphics.mesh.get_object_names(),
        "contour": session.results.graphics.contour.get_object_names(),
        "vector": session.results.graphics.vector.get_object_names(),
        "pathlines": session.results.graphics.pathline.get_object_names(),
        "particletracks": session.results.graphics.particle_track.get_object_names(),
    }

    session.tui.display.set.picture.use_window_resolution("no")

    # Loop through the objects of each key in the Images dictionary,
    # display, save picture, and insert picture into PPTX
    for key, value in Images.items():
        for image in value:
            graph_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(graph_slide_layout)
            title = slide.shapes.title
            title.text = image
            # session.results.graphics.mesh.display(object_name=image)
            session.tui.display.save_picture(image + ".jpg")
            placeholder = slide.placeholders[14]
            pic = placeholder.insert_picture(image + ".jpg")
            adjust_picture_to_fit(pic)

    # Add residual plot.
    session.tui.plot.residuals()
    session.tui.display.save_picture("residuals.jpg")

    graph_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = "Residuals"

    placeholder = slide.placeholders[14]
    pic = placeholder.insert_picture("residuals.jpg")
    adjust_picture_to_fit(pic)

    # Add charts for each individual report plot.

    # Generate a list of all report plots.
    # Using Scheme Eval until Settings API can be used.
    reportdefs_string = session.scheme_eval.exec(
        ('(ti-menu-load-string "solve/report-plots/list ")',)
    )
    reportdefs = reportdefs_string.split()
    reportdefs.pop(0)

    # Loop through all report plots,
    # generate a slide titled with report plot name, including image of plot
    if reportdefs[0] == "invalid":
        print("There are no report definitions.")
    else:
        for report in reportdefs:
            session.tui.solve.report_plots.plot(report)
            session.tui.display.save_picture(report + ".jpg")

            graph_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(graph_slide_layout)
            title = slide.shapes.title
            title.text = report

            placeholder = slide.placeholders[14]
            pic = placeholder.insert_picture(report + ".jpg")
            adjust_picture_to_fit(pic)

    prs.save("report.pptx")
